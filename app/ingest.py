import os
import pdfplumber
import docx
from pptx import Presentation

def extract_text_from_pdf(path : str) -> str:
    text_parts =[]
    with pdfplumber.open(path) as pdf :
        for page in pdf.pages :
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
    return "\n".join(text_parts)

def extract_text_from_docx(path : str) -> str:
    doc = docx.Document(path)
    paragraphs = [para.text for para in doc.paragraphs]
    return "\n".join(paragraphs)

def extract_text_from_pptx(path : str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text_from_txt(path : str) -> str:
    with open(path, 'r', encoding='utf-8') as file:
        return file.read()

def ingest_document(path : str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        return extract_text_from_pdf(path)
    elif ext == '.docx':
        return extract_text_from_docx(path)
    elif ext == '.pptx':
        return extract_text_from_pptx(path)
    elif ext == '.txt':
        return extract_text_from_txt(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")